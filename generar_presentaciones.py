import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import landscape, letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor

slides_data = [
    {
        "title": "Inteligencia Electoral: Encontrando los Votos Ocultos",
        "subtitle": "Estrategia de Micro-Focalización para la Captación de Indecisos y Abstencionistas – Iván Cepeda 2026.",
        "content": "Bienvenidos. Hoy no vamos a hablar de cómo convencer a la extrema derecha, sino de cómo movilizar a la gran masa silenciada de Medellín. Vamos a hablar de inteligencia de datos aplicada a la campaña de Iván Cepeda, utilizando censo milimétrico para encontrar nuestra mayor oportunidad de crecimiento."
    },
    {
        "title": "La Ilusión del Techo Electoral",
        "content": "• La narrativa tradicional dice que 'Medellín es de derecha'.\n• Los datos muestran que el mayor partido es la Abstención y el Voto en Blanco.\n• Crecer robando votos a la derecha es costoso e ineficiente.\n\nEl mapa tradicional nos pinta un Medellín azul, dominado por la derecha. Pero si miramos de cerca, el verdadero ganador histórico en muchas comunas populares es la desconexión. Nuestro techo electoral no lo define la derecha, lo define la abstención."
    },
    {
        "title": "De la Estimación a la Precisión Mesa a Mesa",
        "content": "• Adiós a las estimaciones por comuna.\n• Uso directo del Censo Electoral Divipole 2026.\n• Cruce espacial: Se mapearon 2403 puestos exactos hacia sus respectivos barrios.\n\nPara esta campaña, descartamos las estimaciones a ojo. Integramos el Divipole 2026 para conocer exactamente, puesto por puesto, cuántos hombres y mujeres pueden votar. Si en un puesto hay 12,000 habilitados y solo votaron 4,000, sabemos con coordenadas exactas dónde están esos 8,000 votos faltantes."
    },
    {
        "title": "La Fórmula de la Victoria",
        "content": "• Oportunidad = (Abstención × Afinidad a Cepeda) + Votos Blancos/Nulos + Votos Alternativos.\n• No buscamos a todos los que no votan; buscamos a los que no votan pero piensan como nosotros.\n\nNo tiene sentido ir a buscar abstencionistas en El Poblado, donde la afinidad histórica nos castiga. El 'Voto Oportunidad' pondera matemáticamente la abstención por el nivel de favorabilidad hacia la izquierda en ese barrio específico. Buscamos bolsas de abstención en terrenos fértiles."
    },
    {
        "title": "Geografía de la Indecisión",
        "content": "• Las zonas céntricas y de estrato alto tienen alta participación y voto duro de derecha.\n• Las zonas periféricas (Nororiente, Noroccidente, Corregimientos) concentran el potencial inactivo.\n\nCuando superponemos nuestra fórmula en el mapa, la ciudad cambia. Las zonas rojas y moradas nos muestran dónde el estado y las campañas tradicionales no han llegado. Aquí es donde Iván Cepeda tiene margen de crecimiento a bajo costo."
    },
    {
        "title": "Micro-Focalización: El Top 50",
        "content": "• Identificamos los 50 barrios con el mayor 'Voto Oportunidad'.\n• Representan el ROI (Retorno de Inversión) político más alto.\n• Concentración de esfuerzos de campaña (Tierra y Aire).\n\nEn una campaña, el tiempo y los recursos son finitos. Hemos aislado el Top 50 de barrios críticos. Si concentramos nuestra logística, pauta y recorridos físicos exclusivamente en estos sectores, maximizamos cada peso y cada hora invertida."
    },
    {
        "title": "¿A Quién le Estamos Hablando?",
        "content": "• Desencantados: Sienten que 'todos son iguales'.\n• Votantes de Opinión Blanda: Votaron en blanco o nulo por falta de conexión emocional.\n• Jóvenes y Mujeres: Segmentos con mayor volatilidad en nuestro Divipole.\n\nEste potencial no es militante. No responden a discursos ideológicos densos. Son personas que sufren los problemas reales de la ciudad y sienten que la política los abandonó. Nuestro mensaje debe ser pragmático, empático y directo."
    },
    {
        "title": "El Mensaje: Hechos, No Promesas",
        "content": "• Contraste: La derecha habla de miedo; Cepeda debe hablar de paz material (seguridad económica, justicia social).\n• Trayectoria: Validar la lucha anticorrupción y por la paz de Cepeda como garantía de cumplimiento.\n• Llamado a la acción: 'Tu silencio le da el poder a los mismos de siempre.'\n\nA un indeciso no se le convence con un debate de izquierda vs. derecha. Se le convence demostrando que su voto sí cambia su realidad. La narrativa de Iván Cepeda debe enfocarse en su coherencia histórica: el hombre que enfrenta a los poderosos, ahora viene a defender tu barrio."
    },
    {
        "title": "Táctica Territorial en el Top 50",
        "content": "• Tomas Barriales: Eventos en los puestos de votación de mayor potencial.\n• Comités de Cuidado: Líderes de cuadra encargados de pedagogía electoral.\n• Movilización Focalizada: Activación de mujeres y jóvenes.\n\nCon la aplicación, sabemos exactamente a qué cuadra y a qué escuela ir. El trabajo en tierra será de pedagogía uno-a-uno en el Top 50. Activaremos a líderes locales no para convencer al convencido, sino para enseñar a votar al que nunca lo hace."
    },
    {
        "title": "Estrategia Digital Geo-Focalizada",
        "content": "• Anuncios digitales (Meta, YouTube) segmentados exclusivamente en los polígonos del Top 50.\n• Mensajes hiper-locales: 'En el barrio [X], merecemos más'.\n• Retargeting a usuarios que interactúan con contenido de paz o justicia social.\n\nNo gastaremos dinero mostrando anuncios en sectores donde ya sabemos que no nos van a votar. Toda nuestra pauta digital estará delimitada por GPS a los polígonos de nuestro Top 50. Será una campaña invisible para el resto de la ciudad, pero omnipresente para nuestros indecisos."
    },
    {
        "title": "Control y Seguimiento",
        "content": "• KPI 1: Reducción del porcentaje de indecisos en encuestas focalizadas en el Top 50.\n• KPI 2: Tasa de reclutamiento de testigos electorales en los puestos de alta 'Oportunidad'.\n• KPI 3: Interacción digital y CTR (Click-through rate) por comuna.\n\nLa estrategia basada en datos requiere medición constante. Monitorearemos nuestro impacto semanalmente, no a nivel ciudad, sino barrio por barrio dentro de nuestra lista prioritaria. Ajustaremos la pauta y las visitas en tiempo real."
    },
    {
        "title": "El Camino a la Victoria",
        "content": "• Los votos existen; solo están inactivos.\n• Tenemos la tecnología (App, Divipole, Georreferenciación) para encontrarlos.\n• Iván Cepeda representa el perfil exacto de coherencia para activar este voto indignado.\n\nPara ganar no necesitamos milagros, necesitamos método. La ciencia de datos nos ha dado el mapa del tesoro. Ahora, con el mensaje correcto y la presencia de Iván Cepeda en el territorio exacto, vamos a convertir la abstención en nuestra mayor fortaleza electoral. Muchas gracias."
    }
]

def generate_html():
    html = '''<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Estrategia Iván Cepeda 2026</title>
    <style>
        body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; margin: 0; padding: 0; background-color: #f8fafc; color: #0f172a; }
        .slide { min-height: 100vh; display: flex; flex-direction: column; justify-content: center; align-items: center; padding: 40px; box-sizing: border-box; border-bottom: 2px solid #e2e8f0; page-break-after: always; }
        .slide-content { max-width: 900px; width: 100%; background: white; padding: 50px; border-radius: 12px; box-shadow: 0 10px 25px rgba(0,0,0,0.05); }
        h1 { color: #ea580c; font-size: 2.5rem; margin-top: 0; }
        h2 { color: #ea580c; font-size: 2rem; margin-top: 0; }
        p, ul { font-size: 1.25rem; line-height: 1.6; color: #334155; }
        ul { padding-left: 30px; }
        li { margin-bottom: 10px; }
        .notes { margin-top: 30px; padding: 20px; background-color: #f1f5f9; border-left: 4px solid #3b82f6; border-radius: 4px; font-style: italic; font-size: 1.1rem; }
        .slide-number { text-align: center; margin-top: 20px; color: #94a3b8; font-size: 0.9rem; }
        @media print {
            body { background-color: white; }
            .slide { min-height: auto; padding: 0; border: none; box-shadow: none; display: block; margin-bottom: 50px; }
            .slide-content { box-shadow: none; padding: 20px; }
        }
    </style>
</head>
<body>
'''
    for i, slide in enumerate(slides_data):
        html += f'    <div class="slide" id="slide-{i+1}">\n'
        html += f'        <div class="slide-content">\n'
        if i == 0:
            html += f'            <h1>{slide["title"]}</h1>\n'
            html += f'            <p style="font-size: 1.5rem; font-weight: bold;">{slide["subtitle"]}</p>\n'
            html += f'            <div class="notes"><strong>Notas del Orador:</strong><br>{slide["content"]}</div>\n'
        else:
            html += f'            <h2>{slide["title"]}</h2>\n'
            
            content_parts = slide["content"].split('\n\n')
            bullets = content_parts[0].split('\n')
            
            html += f'            <ul>\n'
            for b in bullets:
                if b.startswith('• '): b = b[2:]
                html += f'                <li>{b}</li>\n'
            html += f'            </ul>\n'
            
            if len(content_parts) > 1:
                html += f'            <div class="notes"><strong>Notas del Orador:</strong><br>{content_parts[1]}</div>\n'
                
        html += f'            <div class="slide-number">Diapositiva {i+1} de {len(slides_data)}</div>\n'
        html += f'        </div>\n'
        html += f'    </div>\n'

    html += '''
</body>
</html>
'''
    with open('Presentacion_Estrategia_Cepeda.html', 'w', encoding='utf-8') as f:
        f.write(html)
        
def generate_pdf():
    doc = SimpleDocTemplate("Presentacion_Estrategia_Cepeda.pdf", pagesize=landscape(letter))
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle('TitleStyle', parent=styles['Heading1'], fontSize=28, textColor=HexColor('#ea580c'), alignment=1, spaceAfter=20)
    subtitle_style = ParagraphStyle('SubStyle', parent=styles['Normal'], fontSize=18, textColor=HexColor('#334155'), alignment=1, spaceAfter=40)
    
    heading_style = ParagraphStyle('HeadingStyle', parent=styles['Heading1'], fontSize=24, textColor=HexColor('#ea580c'), spaceAfter=20)
    bullet_style = ParagraphStyle('BulletStyle', parent=styles['Normal'], fontSize=16, leading=22, spaceAfter=10, leftIndent=20)
    notes_style = ParagraphStyle('NotesStyle', parent=styles['Normal'], fontSize=14, leading=18, textColor=HexColor('#475569'), fontName='Helvetica-Oblique', spaceBefore=20)
    
    story = []
    
    for i, slide in enumerate(slides_data):
        if i == 0:
            story.append(Spacer(1, 100))
            story.append(Paragraph(slide["title"], title_style))
            story.append(Paragraph(slide["subtitle"], subtitle_style))
            story.append(Spacer(1, 40))
            story.append(Paragraph("<b>Notas del Orador:</b> " + slide["content"], notes_style))
            story.append(PageBreak())
        else:
            story.append(Spacer(1, 40))
            story.append(Paragraph(slide["title"], heading_style))
            
            content_parts = slide["content"].split('\n\n')
            bullets = content_parts[0].split('\n')
            
            for b in bullets:
                b = b.replace('• ', '')
                story.append(Paragraph("• " + b, bullet_style))
                
            if len(content_parts) > 1:
                story.append(Spacer(1, 20))
                story.append(Paragraph("<b>Notas del Orador:</b> " + content_parts[1], notes_style))
            
            story.append(PageBreak())
            
    doc.build(story)

def generate_pptx():
    prs = Presentation()
    
    for i, slide in enumerate(slides_data):
        if i == 0:
            slide_layout = prs.slide_layouts[0] # Title slide
            s = prs.slides.add_slide(slide_layout)
            title = s.shapes.title
            subtitle = s.placeholders[1]
            title.text = slide["title"]
            subtitle.text = slide["subtitle"]
            
            # Add notes
            notes_slide = s.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide["content"]
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content
            s = prs.slides.add_slide(slide_layout)
            title = s.shapes.title
            content_box = s.placeholders[1]
            
            title.text = slide["title"]
            
            content_parts = slide["content"].split('\n\n')
            bullets = content_parts[0].split('\n')
            
            tf = content_box.text_frame
            for j, b in enumerate(bullets):
                b = b.replace('• ', '')
                if j == 0:
                    tf.text = b
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
            
            # Add notes
            if len(content_parts) > 1:
                notes_slide = s.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = content_parts[1]

    prs.save('Presentacion_Estrategia_Cepeda.pptx')

print("Generando HTML...")
generate_html()
print("Generando PDF...")
generate_pdf()
print("Generando PPTX...")
generate_pptx()
print("¡Archivos generados exitosamente!")
